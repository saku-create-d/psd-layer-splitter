import streamlit as st
import zipfile
import io
import re
import math
from PIL import Image, ImageDraw, ImageFont
from psd_tools import PSDImage
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ── ページ設定 ────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="File Layer Splitter",
    page_icon="🎨",
    layout="centered",
)

# ── カスタム CSS ───────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Space+Mono:wght@400;700&family=Syne:wght@700;800&display=swap');

html, body, [class*="css"] {
    font-family: 'Space Mono', monospace;
}

.stApp {
    background: #0a0a0f;
    background-image:
        radial-gradient(ellipse 80% 50% at 20% 20%, rgba(100, 60, 255, 0.12) 0%, transparent 60%),
        radial-gradient(ellipse 60% 40% at 80% 80%, rgba(255, 60, 120, 0.08) 0%, transparent 60%);
}

.hero {
    text-align: center;
    padding: 3rem 0 2rem;
    border-bottom: 1px solid rgba(255,255,255,0.06);
    margin-bottom: 2.5rem;
}
.hero h1 {
    font-family: Helvetica, Arial, sans-serif;
    font-weight: 800;
    font-size: 3rem;
    letter-spacing: -0.02em;
    background: linear-gradient(135deg, #ffb347 0%, #ff8c00 50%, #ff6a00 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
    margin: 0 0 0.4rem;
    line-height: 1.1;
}
.hero p {
    color: rgba(255,255,255,0.38);
    font-size: 0.78rem;
    letter-spacing: 0.18em;
    text-transform: uppercase;
    margin: 0;
}

.badge-row {
    display: flex;
    justify-content: center;
    gap: 8px;
    flex-wrap: wrap;
    margin-top: 1rem;
}
.badge {
    font-family: 'Space Mono', monospace;
    font-size: 0.65rem;
    font-weight: 700;
    letter-spacing: 0.1em;
    padding: 3px 10px;
    border-radius: 20px;
    border: 1px solid;
}
.badge-psd  { color: #b388ff; border-color: rgba(179,136,255,0.4); background: rgba(179,136,255,0.07); }
.badge-ai   { color: #ff9e40; border-color: rgba(255,158,64,0.4);  background: rgba(255,158,64,0.07);  }
.badge-eps  { color: #40c8ff; border-color: rgba(64,200,255,0.4);  background: rgba(64,200,255,0.07);  }
.badge-pdf  { color: #ff6060; border-color: rgba(255,96,96,0.4);   background: rgba(255,96,96,0.07);   }
.badge-jpg  { color: #60ff9e; border-color: rgba(96,255,158,0.4);  background: rgba(96,255,158,0.07);  }
.badge-pptx { color: #ffdd57; border-color: rgba(255,221,87,0.4);  background: rgba(255,221,87,0.07);  }

[data-testid="stFileUploader"] {
    background: rgba(255,255,255,0.03) !important;
    border: 1px dashed rgba(179, 136, 255, 0.35) !important;
    border-radius: 12px !important;
    padding: 1.5rem !important;
    transition: border-color 0.2s ease;
}
[data-testid="stFileUploader"]:hover {
    border-color: rgba(179, 136, 255, 0.65) !important;
}

.stButton > button,
.stDownloadButton > button {
    font-family: 'Space Mono', monospace !important;
    font-weight: 700 !important;
    font-size: 0.78rem !important;
    letter-spacing: 0.12em !important;
    text-transform: uppercase !important;
    border-radius: 6px !important;
    padding: 0.65rem 1.8rem !important;
    transition: all 0.2s ease !important;
}
.stDownloadButton > button {
    background: linear-gradient(135deg, #7c4dff, #ff4081) !important;
    border: none !important;
    color: #fff !important;
    width: 100% !important;
}
.stDownloadButton > button:hover {
    opacity: 0.88 !important;
    transform: translateY(-1px) !important;
    box-shadow: 0 8px 24px rgba(124, 77, 255, 0.4) !important;
}

[data-testid="stMetric"] {
    background: rgba(255,255,255,0.03);
    border: 1px solid rgba(255,255,255,0.07);
    border-radius: 10px;
    padding: 1rem 1.2rem !important;
}
[data-testid="stMetricLabel"] { color: rgba(255,255,255,0.35) !important; font-size: 0.68rem !important; }
[data-testid="stMetricValue"] { color: #b388ff !important; font-family: 'Syne', sans-serif !important; }

hr { border-color: rgba(255,255,255,0.06) !important; }
p, li { color: rgba(255,255,255,0.6) !important; }
h2, h3 { color: rgba(255,255,255,0.85) !important; font-family: 'Syne', sans-serif !important; }

[data-testid="stProgress"] > div > div {
    background: linear-gradient(90deg, #7c4dff, #ff4081) !important;
    border-radius: 4px !important;
}
[data-testid="stAlert"] {
    border-radius: 10px !important;
    font-size: 0.82rem !important;
}

::-webkit-scrollbar { width: 4px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: rgba(179,136,255,0.3); border-radius: 4px; }
</style>
""", unsafe_allow_html=True)

# ── ヘッダー ──────────────────────────────────────────────────────────────────
st.markdown("""
<div class="hero">
    <h1>File Layer Splitter</h1>
    <p>Upload · Extract · Download</p>
    <div class="badge-row">
        <span class="badge badge-psd">PSD</span>
        <span class="badge badge-ai">AI</span>
        <span class="badge badge-eps">EPS</span>
        <span class="badge badge-pdf">PDF</span>
        <span class="badge badge-jpg">JPG</span>
        <span class="badge badge-pptx">PPTX</span>
    </div>
</div>
""", unsafe_allow_html=True)


# ════════════════════════════════════════════════════════════════════════════════
# 共通ヘルパー
# ════════════════════════════════════════════════════════════════════════════════

def sanitize_name(name: str) -> str:
    """ファイル名として安全な文字列に変換"""
    return re.sub(r'[\\/:*?"<>|]', "_", name).strip() or "shape"


def pil_to_png_bytes(img: Image.Image) -> bytes:
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def make_thumb(img: Image.Image, size: int = 120) -> bytes:
    """サムネイル用バイト列（暗色背景に RGBA 合成）"""
    thumb = img.copy()
    thumb.thumbnail((size, size))
    if thumb.mode == "RGBA":
        bg = Image.new("RGB", thumb.size, (30, 30, 30))
        bg.paste(thumb, mask=thumb.split()[3])
        thumb = bg
    elif thumb.mode != "RGB":
        thumb = thumb.convert("RGB")
    return pil_to_png_bytes(thumb)


# ════════════════════════════════════════════════════════════════════════════════
# PSD 処理
# ════════════════════════════════════════════════════════════════════════════════

def collect_layers(layer, depth=0):
    results = []
    if isinstance(layer, PSDImage):
        for child in layer:
            results.extend(collect_layers(child, depth))
    elif layer.is_group():
        for child in layer:
            results.extend(collect_layers(child, depth + 1))
    else:
        results.append(layer)
    return results


def process_psd(file_bytes: bytes):
    psd = PSDImage.open(io.BytesIO(file_bytes))
    layers = collect_layers(psd)
    results = []
    name_count: dict[str, int] = {}

    for layer in layers:
        try:
            img = layer.composite()
        except Exception:
            continue
        if img is None:
            continue

        base = sanitize_name(layer.name)
        if base in name_count:
            name_count[base] += 1
            fname = f"{base}_{name_count[base]:02d}.png"
        else:
            name_count[base] = 0
            fname = f"{base}.png"

        results.append((fname, layer.name, pil_to_png_bytes(img), img))

    meta = {"width": psd.width, "height": psd.height}
    return results, meta


# ════════════════════════════════════════════════════════════════════════════════
# AI / EPS / PDF 処理（PyMuPDF）
# ════════════════════════════════════════════════════════════════════════════════

def _fitz_filetype(name: str) -> str:
    ext = name.rsplit(".", 1)[-1].lower()
    return {"ai": "pdf", "eps": "eps", "pdf": "pdf"}.get(ext, "pdf")


def process_via_fitz(file_bytes: bytes, original_name: str):
    doc = fitz.open(stream=file_bytes, filetype=_fitz_filetype(original_name))
    results = []
    matrix = fitz.Matrix(2, 2)

    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=matrix, alpha=True)
        png_bytes = pix.tobytes("png")
        img = Image.open(io.BytesIO(png_bytes))
        results.append((f"page_{i + 1:03d}.png", f"Page {i + 1}", png_bytes, img))

    meta = {
        "pages":  doc.page_count,
        "width":  int(doc[0].rect.width * 2)  if doc.page_count else 0,
        "height": int(doc[0].rect.height * 2) if doc.page_count else 0,
    }
    doc.close()
    return results, meta


# ════════════════════════════════════════════════════════════════════════════════
# JPG / JPEG 処理
# ════════════════════════════════════════════════════════════════════════════════

def process_jpg(file_bytes: bytes, original_name: str):
    img = Image.open(io.BytesIO(file_bytes)).convert("RGBA")
    base = sanitize_name(original_name.rsplit(".", 1)[0])
    meta = {"width": img.width, "height": img.height}
    return [(f"{base}.png", original_name, pil_to_png_bytes(img), img)], meta


# ════════════════════════════════════════════════════════════════════════════════
# PPTX 処理：スライド内シェイプを個別 PNG 化
# ════════════════════════════════════════════════════════════════════════════════

# EMU → ピクセル変換（96 dpi 相当を 2× = 192 dpi で出力）
_DPI_SCALE = 2
_EMU_PER_INCH = 914400
_PX_PER_INCH = 96 * _DPI_SCALE  # 192 px/inch


def _emu_to_px(emu: int) -> int:
    return max(1, round(emu / _EMU_PER_INCH * _PX_PER_INCH))


def _shape_type_label(shape) -> str:
    """シェイプ種別を人間が読める文字列で返す"""
    try:
        st_val = shape.shape_type
        if st_val == MSO_SHAPE_TYPE.PICTURE:
            return "Picture"
        if st_val == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return "AutoShape"
        if st_val == MSO_SHAPE_TYPE.TEXT_BOX:
            return "TextBox"
        if st_val == MSO_SHAPE_TYPE.GROUP:
            return "Group"
        if st_val == MSO_SHAPE_TYPE.TABLE:
            return "Table"
        if st_val == MSO_SHAPE_TYPE.CHART:
            return "Chart"
        if st_val == MSO_SHAPE_TYPE.FREEFORM:
            return "Freeform"
        if st_val == MSO_SHAPE_TYPE.LINE:
            return "Line"
        if st_val == MSO_SHAPE_TYPE.PLACEHOLDER:
            return "Placeholder"
        return f"Type{int(st_val)}"
    except Exception:
        return "Unknown"


def _extract_picture_shape(shape) -> Image.Image | None:
    """
    Picture シェイプから画像バイト列を直接取得して PIL Image で返す。
    失敗した場合は None。
    """
    try:
        image_bytes = shape.image.blob
        return Image.open(io.BytesIO(image_bytes)).convert("RGBA")
    except Exception:
        return None


def _render_non_picture_shape(shape, slide_width_emu: int, slide_height_emu: int) -> Image.Image | None:
    """
    Picture 以外のシェイプを PyMuPDF 経由でスライド全体からクリップして返す。
    シェイプの座標・サイズが有効な場合のみ実行する。

    処理フロー:
      1. shape の left/top/width/height を確認（None なら諦める）
      2. shape が存在するスライドを含む一時 PPTX を python-pptx で構築し
         バイト列に書き出す（直接操作は不可なため）
      3. PyMuPDF でその PPTX を PDF としてレンダリング（※ python-pptx には
         レンダリング機能がないため fitz で代替）
      4. シェイプの PPTX 座標を PDF 座標系に変換してクリップ
    """
    try:
        left   = shape.left
        top    = shape.top
        width  = shape.width
        height = shape.height
    except Exception:
        return None

    if None in (left, top, width, height) or width <= 0 or height <= 0:
        return None

    # ── スライド全体を fitz で PNG レンダリング ────────────────────────────
    # PPTX を一時バイト列として PyMuPDF に渡す
    # （shape._element の親ツリーから prs を復元するのは困難なため、
    #   呼び出し元から slide_pptx_bytes を受け取る設計にする）
    # この関数は _render_shape_via_slide_crop() から呼ばれる
    return None  # スタブ — 実際の処理は _render_shape_via_slide_crop で行う


def _render_shape_via_slide_crop(
    shape,
    slide_png: Image.Image,
    slide_width_emu: int,
    slide_height_emu: int,
) -> Image.Image | None:
    """
    既にレンダリング済みのスライド画像から、shape の座標領域を切り抜く。
    """
    try:
        left   = shape.left   or 0
        top    = shape.top    or 0
        width  = shape.width  or 0
        height = shape.height or 0
    except Exception:
        return None

    if width <= 0 or height <= 0:
        return None

    sw, sh = slide_png.size

    # EMU → スライド画像ピクセル座標に変換
    sx = round(left   / slide_width_emu  * sw)
    sy = round(top    / slide_height_emu * sh)
    ex = round((left + width)  / slide_width_emu  * sw)
    ey = round((top  + height) / slide_height_emu * sh)

    # 画像境界でクランプ
    sx, sy = max(0, sx), max(0, sy)
    ex, ey = min(sw, ex), min(sh, ey)

    if ex <= sx or ey <= sy:
        return None

    return slide_png.crop((sx, sy, ex, ey))


def _slide_to_pil(prs_bytes: bytes, slide_index: int) -> Image.Image | None:
    """
    PPTX バイト列の指定スライドを PyMuPDF でレンダリングして PIL Image を返す。
    PyMuPDF は PPTX を直接レンダリングできないため、
    まず python-pptx で該当スライドのみの一時 PPTX を作り、
    fitz で PDF として開いてラスタライズする。
    （注: fitz は .pptx を直接開けないため PDF 経由は不可。
      代わりに .pptx を zip として扱い内部の slide.xml を SVG/PNG に
      変換する手段も存在するが、依存が増えるため、ここでは
      オリジナル PPTX 全体を fitz に渡して目的のページだけ使う。）
    """
    try:
        # fitz は .pptx を直接開けないため、
        # python-pptx でスライド画像を生成する代わりに
        # PPTX をそのまま fitz に渡す（fitz >= 1.22 では pptx 対応あり）
        doc = fitz.open(stream=prs_bytes, filetype="pptx")
        if slide_index >= doc.page_count:
            doc.close()
            return None
        matrix = fitz.Matrix(_DPI_SCALE, _DPI_SCALE)
        pix = doc[slide_index].get_pixmap(matrix=matrix, alpha=False)
        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        doc.close()
        return img.convert("RGBA")
    except Exception:
        return None


def process_pptx(file_bytes: bytes):
    """
    PPTX の各スライドを走査し、シェイプを個別 PNG として抽出する。

    戦略（優先順）:
      1. Picture シェイプ → image.blob を直接デコード（最高品質）
      2. その他のシェイプ → スライド全体をレンダリングして座標でクロップ
         - まず PyMuPDF (fitz) で PPTX をレンダリング（fitz >= 1.22 対応時）
         - 失敗した場合は Pillow でプレースホルダー画像を生成

    Returns:
        items : list of (file_name, label, png_bytes, PIL.Image)
        meta  : dict
    """
    prs = Presentation(io.BytesIO(file_bytes))
    results: list[tuple[str, str, bytes, Image.Image]] = []
    total_shapes = 0

    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        # スライド全体のレンダリング（非 Picture シェイプのクロップ用）
        slide_png: Image.Image | None = _slide_to_pil(file_bytes, slide_idx)

        for shape_idx, shape in enumerate(slide.shapes):
            total_shapes += 1
            shape_num = shape_idx + 1
            type_label = _shape_type_label(shape)
            shape_name = sanitize_name(shape.name) if shape.name else f"shape{shape_num}"
            fname = f"slide{slide_num:02d}_{shape_num:03d}_{shape_name}.png"
            label = f"Slide {slide_num} / {shape.name or f'Shape {shape_num}'} [{type_label}]"

            img: Image.Image | None = None

            # ── 戦略 1: Picture シェイプは blob を直接取得 ──────────────────
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = _extract_picture_shape(shape)

            # ── 戦略 2: スライドレンダリング画像からクロップ ─────────────────
            if img is None and slide_png is not None:
                img = _render_shape_via_slide_crop(
                    shape, slide_png, slide_w_emu, slide_h_emu
                )

            # ── 戦略 3: フォールバック — プレースホルダー画像を生成 ───────────
            if img is None:
                try:
                    w_px = _emu_to_px(shape.width  or _EMU_PER_INCH)
                    h_px = _emu_to_px(shape.height or _EMU_PER_INCH // 2)
                except Exception:
                    w_px, h_px = 200, 100

                w_px = max(w_px, 60)
                h_px = max(h_px, 30)

                img = Image.new("RGBA", (w_px, h_px), (50, 50, 70, 220))
                draw = ImageDraw.Draw(img)
                # 外枠
                draw.rectangle([0, 0, w_px - 1, h_px - 1], outline=(150, 150, 200, 200), width=2)
                # ラベルテキスト
                text = f"[{type_label}]\n{shape.name or ''}"
                try:
                    draw.text((8, 8), text, fill=(200, 200, 220, 255))
                except Exception:
                    pass

            png_bytes = pil_to_png_bytes(img)
            results.append((fname, label, png_bytes, img))

    meta = {
        "slides":       len(prs.slides),
        "total_shapes": total_shapes,
        "width":        _emu_to_px(slide_w_emu),
        "height":       _emu_to_px(slide_h_emu),
    }
    return results, meta


# ════════════════════════════════════════════════════════════════════════════════
# PDF 要素抽出処理（画像 + ベクター図形）
# ════════════════════════════════════════════════════════════════════════════════

# ページレンダリング解像度（図形クロップ用）
_PDF_MATRIX = fitz.Matrix(2, 2)  # 2× ≈ 144 dpi


def _rect_valid(rect: fitz.Rect, min_px: float = 4.0) -> bool:
    """クロップ対象として有効なサイズかチェック（2×スケール後）"""
    w = (rect.x1 - rect.x0) * 2
    h = (rect.y1 - rect.y0) * 2
    return w >= min_px and h >= min_px


def _extract_images_from_page(
    doc: fitz.Document,
    page: fitz.Page,
    page_num: int,
    zip_folder: str,
) -> list[tuple[str, str, bytes, "Image.Image"]]:
    """
    page.get_images() で埋め込み画像を抽出する。
    同一 xref は重複スキップ。

    Returns: list of (zip_path, label, png_bytes, PIL.Image)
    """
    results = []
    seen_xrefs: set[int] = set()

    for img_idx, img_info in enumerate(page.get_images(full=True)):
        xref = img_info[0]
        if xref in seen_xrefs:
            continue
        seen_xrefs.add(xref)

        try:
            base_img = doc.extract_image(xref)
            raw_bytes = base_img["image"]
            ext_hint  = base_img.get("ext", "png")

            # PIL で開いて RGBA に統一
            pil_img = Image.open(io.BytesIO(raw_bytes)).convert("RGBA")
            png_bytes = pil_to_png_bytes(pil_img)

            fname     = f"img_{img_idx + 1:03d}_xref{xref}.png"
            zip_path  = f"{zip_folder}/{fname}"
            label     = f"P{page_num} / Image {img_idx + 1} (xref={xref}, orig={ext_hint})"
            results.append((zip_path, label, png_bytes, pil_img))

        except Exception:
            continue

    return results


def _cluster_drawings(drawings: list[dict], gap: float = 4.0) -> list[fitz.Rect]:
    """
    get_drawings() の各要素の rect を「近接するもの同士」でクラスタリングし、
    グループごとの合成 Rect リストを返す。
    これにより、同一論理図形を構成する複数パスをまとめて1枚の PNG にできる。

    gap: この距離（PDF pt 単位）以内なら同一グループとみなす
    """
    if not drawings:
        return []

    rects = []
    for d in drawings:
        r = d.get("rect")
        if r and not fitz.Rect(r).is_empty:
            rects.append(fitz.Rect(r))

    if not rects:
        return []

    # Union-Find でクラスタリング
    n = len(rects)
    parent = list(range(n))

    def find(x):
        while parent[x] != x:
            parent[x] = parent[parent[x]]
            x = parent[x]
        return x

    def union(a, b):
        parent[find(a)] = find(b)

    for i in range(n):
        for j in range(i + 1, n):
            ri, rj = rects[i], rects[j]
            # 拡張 Rect が重なるか gap 以内ならマージ
            expanded_i = ri + (-gap, -gap, gap, gap)
            if not expanded_i.intersects(rj):
                continue
            union(i, j)

    # グループごとに Rect を合成
    groups: dict[int, fitz.Rect] = {}
    for i, r in enumerate(rects):
        root = find(i)
        if root not in groups:
            groups[root] = fitz.Rect(r)
        else:
            groups[root] = groups[root] | r

    return list(groups.values())


def _extract_drawings_from_page(
    page: fitz.Page,
    page_img: "Image.Image",  # 2× レンダリング済み PIL Image
    page_num: int,
    zip_folder: str,
    min_area_pt: float = 16.0,
) -> list[tuple[str, str, bytes, "Image.Image"]]:
    """
    page.get_drawings() でベクター図形の範囲を特定し、
    レンダリング済みページ画像からクロップして PNG 化する。

    min_area_pt: この面積（pt²）未満の Rect はスキップ
    """
    results = []
    drawings = page.get_drawings()
    if not drawings:
        return results

    clusters = _cluster_drawings(drawings)
    page_rect = page.rect  # PDF pt 単位のページ全体 Rect
    pw, ph = page_img.size  # レンダリング済み画像のピクセルサイズ

    for idx, cluster_rect in enumerate(clusters):
        # 最小面積フィルタ
        area = (cluster_rect.x1 - cluster_rect.x0) * (cluster_rect.y1 - cluster_rect.y0)
        if area < min_area_pt:
            continue
        if not _rect_valid(cluster_rect):
            continue

        # PDF 座標 → ピクセル座標（2× スケール）
        sx = max(0, round((cluster_rect.x0 - page_rect.x0) / page_rect.width  * pw))
        sy = max(0, round((cluster_rect.y0 - page_rect.y0) / page_rect.height * ph))
        ex = min(pw, round((cluster_rect.x1 - page_rect.x0) / page_rect.width  * pw))
        ey = min(ph, round((cluster_rect.y1 - page_rect.y0) / page_rect.height * ph))

        if ex <= sx or ey <= sy:
            continue

        cropped   = page_img.crop((sx, sy, ex, ey)).convert("RGBA")
        png_bytes = pil_to_png_bytes(cropped)

        fname    = f"drawing_{idx + 1:03d}.png"
        zip_path = f"{zip_folder}/{fname}"
        label    = (
            f"P{page_num} / Drawing {idx + 1} "
            f"({cluster_rect.x0:.0f},{cluster_rect.y0:.0f}–"
            f"{cluster_rect.x1:.0f},{cluster_rect.y1:.0f} pt)"
        )
        results.append((zip_path, label, png_bytes, cropped))

    return results


def process_pdf_elements(file_bytes: bytes, original_name: str):
    """
    PDF の各ページから Image と Drawing を個別に抽出する。

    ZIP 構成:
        {pdf_stem}/
            page_{N:03d}/
                img_{M:03d}_xref{X}.png     ← 埋め込み画像
                drawing_{K:03d}.png         ← ベクター図形クロップ

    Returns:
        items : list of (zip_path, label, png_bytes, PIL.Image)
        meta  : dict
    """
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pdf_stem = sanitize_name(original_name.rsplit(".", 1)[0])
    items: list[tuple[str, str, bytes, Image.Image]] = []

    for page_idx in range(doc.page_count):
        page     = doc[page_idx]
        page_num = page_idx + 1
        folder   = f"{pdf_stem}/page_{page_num:03d}"

        # ── ① 埋め込み画像の抽出 ───────────────────────────────────────────
        img_items = _extract_images_from_page(doc, page, page_num, folder)
        items.extend(img_items)

        # ── ② ページ全体をレンダリング（Drawing クロップ用） ──────────────
        pix      = page.get_pixmap(matrix=_PDF_MATRIX, alpha=False)
        page_pil = Image.frombytes("RGB", (pix.width, pix.height), pix.samples).convert("RGBA")

        # ── ③ ベクター図形の抽出 ─────────────────────────────────────────
        drw_items = _extract_drawings_from_page(page, page_pil, page_num, folder)
        items.extend(drw_items)

    meta = {
        "pages":   doc.page_count,
        "images":  sum(1 for it in items if "/img_"     in it[0]),
        "drawings":sum(1 for it in items if "/drawing_" in it[0]),
    }
    doc.close()
    return items, meta


# ════════════════════════════════════════════════════════════════════════════════
# ZIP 生成
# ════════════════════════════════════════════════════════════════════════════════

def build_zip(items: list) -> bytes:
    """
    items の各エントリ fname はそのまま ZIP 内パスとして使う。
    階層付き（例: "my.pdf/page01/img_001.png"）でも平坦でも対応。
    """
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for fname, _, png_bytes, _ in items:
            zf.writestr(fname, png_bytes)
    return buf.getvalue()


# ════════════════════════════════════════════════════════════════════════════════
# メイン UI
# ════════════════════════════════════════════════════════════════════════════════

uploaded = st.file_uploader(
    "ファイルをドロップ、またはクリックして選択",
    type=["psd", "ai", "eps", "pdf", "jpg", "jpeg", "pptx"],
    help="対応形式: PSD / AI / EPS / PDF / JPG / PPTX",
)

if uploaded is not None:
    st.divider()
    ext = uploaded.name.rsplit(".", 1)[-1].lower()
    file_bytes = uploaded.read()

    spinner_msg = {
        "psd":  "PSD レイヤーを解析中...",
        "ai":   "AI ファイルをレンダリング中...",
        "eps":  "EPS ファイルをレンダリング中...",
        "pdf":  "PDF 要素（画像・図形）を抽出中...",
        "jpg":  "JPG を変換中...",
        "jpeg": "JPG を変換中...",
        "pptx": "PPTX シェイプを抽出中...",
    }.get(ext, "処理中...")

    with st.spinner(spinner_msg):
        try:
            if ext == "psd":
                items, meta = process_psd(file_bytes)
            elif ext in ("ai", "eps"):
                items, meta = process_via_fitz(file_bytes, uploaded.name)
            elif ext == "pdf":
                items, meta = process_pdf_elements(file_bytes, uploaded.name)
            elif ext in ("jpg", "jpeg"):
                items, meta = process_jpg(file_bytes, uploaded.name)
            elif ext == "pptx":
                items, meta = process_pptx(file_bytes)
            else:
                st.error(f"未対応の形式です: .{ext}")
                st.stop()
        except Exception as e:
            st.error(f"ファイルの処理に失敗しました: {e}")
            st.stop()

    if not items:
        st.warning("変換可能なコンテンツが見つかりませんでした。")
        st.stop()

    # ── メトリクス ────────────────────────────────────────────────────────────
    col1, col2, col3 = st.columns(3)
    col1.metric("形式", f".{ext.upper()}")

    if ext == "psd":
        col2.metric("レイヤー数", len(items))
        col3.metric("キャンバス", f"{meta['width']} × {meta['height']}")
    elif ext in ("ai", "eps"):
        col2.metric("ページ数", meta.get("pages", len(items)))
        col3.metric("解像度倍率", "2×（≈144dpi）")
    elif ext == "pdf":
        col2.metric("埋め込み画像", meta.get("images", 0))
        col3.metric("ベクター図形", meta.get("drawings", 0))
    elif ext == "pptx":
        col2.metric("シェイプ数", len(items))
        col3.metric("スライド数", meta.get("slides", "–"))
    else:
        col2.metric("画像数", len(items))
        col3.metric("サイズ", f"{meta['width']} × {meta['height']}")

    st.divider()

    # ── プレビュー ────────────────────────────────────────────────────────────
    st.markdown("### プレビュー")
    progress = st.progress(0, text="サムネイルを生成中...")

    num_cols = min(4, len(items))
    cols = st.columns(num_cols)
    for i, (fname, label, png_bytes, img) in enumerate(items):
        progress.progress((i + 1) / len(items), text=f"{fname} を処理中...")
        with cols[i % num_cols]:
            st.image(make_thumb(img), caption=label, use_container_width=True)

    progress.empty()

    # ── ZIP ダウンロード ───────────────────────────────────────────────────────
    st.divider()
    st.markdown("### ダウンロード")

    zip_bytes = build_zip(items)
    zip_name  = f"{uploaded.name.rsplit('.', 1)[0]}_export.zip"
    total_kb  = sum(len(it[2]) for it in items) / 1024

    st.success(f"✓ {len(items)} ファイルを ZIP にパッケージ化しました（合計 {total_kb:.1f} KB）")

    st.download_button(
        label="⬇  ZIP をダウンロード",
        data=zip_bytes,
        file_name=zip_name,
        mime="application/zip",
    )

    with st.expander("含まれるファイル一覧"):
        for fname, label, png_bytes, _ in items:
            st.markdown(
                f"- `{fname}` &nbsp;&nbsp;·&nbsp;&nbsp; "
                f"<span style='color:rgba(255,255,255,0.35)'>{label}</span> "
                f"<span style='color:rgba(255,255,255,0.2);float:right'>{len(png_bytes)//1024} KB</span>",
                unsafe_allow_html=True,
            )

else:
    st.markdown("""
    <div style="
        background: rgba(255,255,255,0.025);
        border: 1px solid rgba(255,255,255,0.06);
        border-radius: 12px;
        padding: 1.6rem 2rem;
        margin-top: 1.5rem;
        font-size: 0.78rem;
        color: rgba(255,255,255,0.35);
        line-height: 2.2;
    ">
    <strong style="color:rgba(255,255,255,0.55)">使い方</strong><br>
    1. ファイルをアップロードエリアにドロップ<br>
    2. 形式に応じて自動処理（レイヤー分割 / ページ変換 / シェイプ抽出 / PNG化）<br>
    3. ZIP ダウンロードボタンで一括取得<br><br>
    <strong style="color:rgba(255,255,255,0.55)">形式別の処理</strong><br>
    <code>.psd</code> &nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;— レイヤーごとに PNG 分割（psd-tools）<br>
    <code>.ai / .eps</code> &nbsp;&nbsp;&nbsp;&nbsp;&nbsp;— 各ページを高解像度 PNG 化（PyMuPDF）<br>
    <code>.pdf</code> &nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;— ページ内の画像・ベクター図形を個別 PNG 抽出（PyMuPDF）<br>
    <code>.jpg / .jpeg</code> &nbsp;&nbsp;&nbsp;— PNG に変換して ZIP 化<br>
    <code>.pptx</code> &nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;— スライド内の各シェイプを個別 PNG 化（python-pptx + PyMuPDF）
    </div>
    """, unsafe_allow_html=True)
